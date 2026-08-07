# src/canvas_api_mcp/extract.py
"""Turn downloaded course files into plain text.

Only formats a student would actually read are supported. Anything else
fails loudly rather than returning bytes the model cannot use.
"""

from __future__ import annotations

import io


class UnsupportedFileType(Exception):
    """Raised when a file's type has no text extractor."""


PDF_TYPES = {"application/pdf"}
DOCX_TYPES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
}
PPTX_TYPES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
}


def _pdf(content: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    return "\n\n".join((page.extract_text() or "") for page in reader.pages).strip()


def _docx(content: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs if p.text).strip()


def _pptx(content: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    chunks: list[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        lines = [
            shape.text
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text
        ]
        if lines:
            chunks.append(f"--- Slide {index} ---\n" + "\n".join(lines))
    return "\n\n".join(chunks).strip()


def extract_text(content: bytes, content_type: str, filename: str) -> str:
    ctype = (content_type or "").split(";")[0].strip().lower()
    lower_name = filename.lower()

    if ctype in PDF_TYPES or lower_name.endswith(".pdf"):
        return _pdf(content)
    if ctype in DOCX_TYPES or lower_name.endswith(".docx"):
        return _docx(content)
    if ctype in PPTX_TYPES or lower_name.endswith(".pptx"):
        return _pptx(content)
    if ctype.startswith("text/") or ctype == "application/json":
        return content.decode("utf-8", errors="replace")

    raise UnsupportedFileType(
        f"No text extractor for {filename!r} (type {content_type!r}). "
        "Supported: PDF, DOCX, PPTX, and plain text."
    )
