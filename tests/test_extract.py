# tests/test_extract.py
import io

import httpx
import pytest
import respx

from canvas_api_mcp.client import CanvasClient
from canvas_api_mcp.config import Config
from canvas_api_mcp.extract import UnsupportedFileType, extract_text
from canvas_api_mcp.tools.content import do_read_file

CFG = Config(base_url="https://canvas.example.edu", token="tok", max_pages=10)
API = "https://canvas.example.edu/api/v1"


def test_plain_text_is_decoded():
    assert extract_text(b"hello world", "text/plain", "a.txt") == "hello world"


def test_docx_paragraphs_are_extracted():
    from docx import Document

    doc = Document()
    doc.add_paragraph("Lecture One")
    doc.add_paragraph("Big-O notation")
    buf = io.BytesIO()
    doc.save(buf)

    text = extract_text(
        buf.getvalue(),
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "lec.docx",
    )
    assert "Lecture One" in text
    assert "Big-O notation" in text


def test_pptx_slide_text_is_extracted():
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Amortised Analysis"
    buf = io.BytesIO()
    prs.save(buf)

    text = extract_text(
        buf.getvalue(),
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "wk3.pptx",
    )
    assert "Amortised Analysis" in text


def test_unsupported_type_raises():
    with pytest.raises(UnsupportedFileType) as exc:
        extract_text(b"\x00\x01", "image/png", "diagram.png")
    assert "png" in str(exc.value).lower() or "image" in str(exc.value).lower()


@respx.mock
async def test_read_file_fetches_metadata_then_content():
    respx.get(f"{API}/files/501").mock(
        return_value=httpx.Response(200, json={
            "id": 501, "display_name": "notes.txt", "content-type": "text/plain",
            "url": "https://files.example.edu/501?verifier=abc",
        })
    )
    respx.get("https://files.example.edu/501").mock(
        return_value=httpx.Response(200, content=b"Kruskal and Prim")
    )
    client = CanvasClient(CFG)
    result = await do_read_file(client, 501)
    await client.aclose()

    assert result["display_name"] == "notes.txt"
    assert result["text"] == "Kruskal and Prim"
    assert result["truncated"] is False


@respx.mock
async def test_read_file_truncates_long_text():
    respx.get(f"{API}/files/502").mock(
        return_value=httpx.Response(200, json={
            "id": 502, "display_name": "big.txt", "content-type": "text/plain",
            "url": "https://files.example.edu/502",
        })
    )
    respx.get("https://files.example.edu/502").mock(
        return_value=httpx.Response(200, content=b"x" * 200)
    )
    client = CanvasClient(CFG)
    result = await do_read_file(client, 502, max_chars=50)
    await client.aclose()

    assert len(result["text"]) == 50
    assert result["truncated"] is True
    assert result["chars"] == 200


@respx.mock
async def test_read_file_reports_metadata_404_as_structured_error():
    """A deleted/inaccessible file_id must come back as the tool's structured
    error contract, not an unhandled CanvasError."""
    respx.get(f"{API}/files/504").mock(
        return_value=httpx.Response(404, json={"errors": [{"message": "not found"}]})
    )
    client = CanvasClient(CFG)
    result = await do_read_file(client, 504)
    await client.aclose()

    assert result["error"] is True
    assert result["status"] == 404


@respx.mock
async def test_read_file_reports_missing_download_url_as_structured_error():
    respx.get(f"{API}/files/505").mock(
        return_value=httpx.Response(200, json={
            "id": 505, "display_name": "ghost.txt", "content-type": "text/plain",
        })
    )
    client = CanvasClient(CFG)
    result = await do_read_file(client, 505)
    await client.aclose()

    assert result["error"] is True
    assert "ghost.txt" in result["message"]


@respx.mock
async def test_read_file_reports_unsupported_type_as_structured_error():
    respx.get(f"{API}/files/503").mock(
        return_value=httpx.Response(200, json={
            "id": 503, "display_name": "photo.png", "content-type": "image/png",
            "url": "https://files.example.edu/503",
        })
    )
    respx.get("https://files.example.edu/503").mock(
        return_value=httpx.Response(200, content=b"\x89PNG")
    )
    client = CanvasClient(CFG)
    result = await do_read_file(client, 503)
    await client.aclose()

    assert result["error"] is True
    assert "photo.png" in result["message"]
